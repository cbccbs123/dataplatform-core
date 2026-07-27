"""
문서(텍스트·PDF·JSON·오피스)를 형식별로 읽어 임베딩 입력용 텍스트 청크 스트림으로 만든다.

per-asset 파이프라인의 텍스트 경로(extract→embed, 그리고 classify 프로브)가 공유하는
'문서 리더 겸 청커'다. 형식별 파서가 원문을 세그먼트(문단/행/페이지/슬라이드 등)로 뽑으면
``_chunk_from_segments`` 가 ``chunk_size`` 안에 그리디로 묶고 ``_apply_overlap`` 이 문자 단위
오버랩을 부여한다. 여기서 LLM 요약·임베딩은 하지 않는다 — 청크를 받아 요약·임베딩하는 쪽은
``src/llm/text_summarizer.py`` · ``src/embedders/text_embedder.py`` 등 소비자다.

지원 형식과 ``file_kind`` 값(``detect_file_kind`` 가 판정해 넘긴다):
- 일반 텍스트: ``file_kind='txt'`` + ``.txt`` (인코딩은 ``choose_encoding`` 으로 추정)
- PDF: ``file_kind='pdf'`` + ``.pdf`` (``pypdf``, 페이지 단위 세그먼트)
- JSON: ``file_kind='json'`` + ``.json`` (txt 와 동일하게 문자 스트림으로 읽음 — 파싱 안 함)
- 오피스: ``file_kind`` 가 ``word``/``excel``/``powerpoint`` 이고 각각 ``.docx``/``.xlsx``/``.pptx``

오피스/PDF 파서(docx·openpyxl·pptx·pypdf)는 해당 분기에서만 지연 import 한다(미사용 형식의
의존성 기동 비용 회피).
"""

from __future__ import annotations

import codecs
from collections.abc import Iterable, Iterator
from pathlib import Path

from src.file.file_type_defs import (
    ALLOWED_TEXT_META_FILE_KINDS,
    OFFICE_FILE_KINDS,
    TEXT_FILE_KINDS,
    MediaKind,
    OfficeKind,
)

_OFFICE_EXTENSIONS: dict[str, frozenset[str]] = {
    OfficeKind.WORD.value: frozenset({".docx"}),
    OfficeKind.EXCEL.value: frozenset({".xlsx"}),
    OfficeKind.POWERPOINT.value: frozenset({".pptx"}),
}
_MAX_INPUT_CHARS = 20_000_000
_OFFICE_FILE_KINDS = OFFICE_FILE_KINDS
_TEXT_FILE_KINDS = TEXT_FILE_KINDS
_ALLOWED_FILE_KINDS = ALLOWED_TEXT_META_FILE_KINDS


def _normalize_file_kind(value: str | None) -> str | None:
    """파일 종류 문자열을 비교 가능한 형태로 통일한다(공백 제거·소문자).

    Returns:
        정규화된 값. 비었으면 ``None``.
    """
    if value is None:
        return None
    s = str(value).strip().lower()
    if s not in _ALLOWED_FILE_KINDS:
        raise ValueError(f"file_kind는 {sorted(_ALLOWED_FILE_KINDS)} 중 하나여야 합니다.")
    return s


def _choose_encoding(path: Path, preferred_encoding: str) -> str:
    """앞부분 샘플을 읽어 **실제로 디코딩되는** 인코딩을 고른다.

    인코딩을 잘못 고르면 글자가 깨진 채로 임베딩까지 흘러가 검색이 조용히 나빠진다. 그래서
    선호 인코딩부터 한국어권에서 흔한 것들까지 차례로 시도해 본다.

    Args:
        path: 대상 파일.
        preferred_encoding: 가장 먼저 시도할 인코딩.

    Returns:
        디코딩에 성공한 인코딩 이름. 전부 실패하면 마지막 후보를 그대로 돌려준다
        (읽기 단계에서 대체 문자로 처리된다).
    """
    # 판정에는 앞 64KiB 면 충분하다 — 파일 전체를 메모리에 올리면 큰 문서에서 낭비가 크다.
    with path.open("rb") as f:
        sample = f.read(65536)
    for enc in (preferred_encoding, "utf-8-sig", "cp949", "euc-kr"):
        try:
            # 2026-07-15 B5: 64KiB 경계가 멀티바이트 문자 **중간**에 걸리면 일반 decode 는 꼬리에서
            # 실패해(정상 utf-8 한글 문서인데) 다음 후보 cp949 가 관대하게 성공 → 오판(모지바케)했다.
            # incremental decoder 의 final=False 는 잘린 꼬리를 "이어질 바이트"로 대기시키고 통과하므로
            # 절단 위치와 무관하게 본문 바이트로만 판정한다(모든 후보 인코딩에 동일 적용·결정적).
            codecs.getincrementaldecoder(enc)().decode(sample, final=False)
            return enc
        except UnicodeDecodeError:
            continue
    return "utf-8"


def _chunk_from_segments(segments: Iterable[str], chunk_size: int) -> Iterator[str]:
    """세그먼트들을 ``chunk_size`` 문자 한도 안에 그리디로 묶어 청크를 흘린다.

    **가능한 한 문단·행·페이지 경계를 지킨다** — 아무 데서나 자르면 문맥이 끊겨 임베딩 품질이
    떨어진다. 다만 한 세그먼트 자체가 한도를 넘으면 어쩔 수 없이 그것만 잘라 낸다.

    Args:
        segments: 문단·행·페이지 단위 조각들. 빈 조각은 버린다.
        chunk_size: 청크 하나의 글자 한도.

    Yields:
        묶인 청크.
    """
    buf = ""
    for seg in segments:
        if not seg:
            continue
        s = seg.strip()
        if not s:
            continue
        if len(s) > chunk_size:
            if buf:
                yield buf
                buf = ""
            for i in range(0, len(s), chunk_size):
                part = s[i : i + chunk_size]
                if part:
                    yield part
            continue
        if not buf:
            buf = s
            continue
        if len(buf) + 1 + len(s) <= chunk_size:
            buf += "\n" + s
        else:
            yield buf
            buf = s
    if buf:
        yield buf


def _limit_segments(segments: Iterable[str], max_chars: int) -> Iterator[str]:
    """세그먼트를 정해진 글자 수까지만 통과시킨다(아주 긴 문서의 상한).

    Args:
        segments: 원본 세그먼트 스트림.
        max_chars: 누적 글자 상한. **0 이하면 아무것도 통과시키지 않는다**.

    Yields:
        상한 안의 세그먼트. 마지막 세그먼트는 상한에 맞춰 **잘라서** 낸다(버리지 않는다).
    """
    if max_chars <= 0:
        return
    remaining = max_chars
    for seg in segments:
        if remaining <= 0:
            break
        if not seg:
            continue
        if len(seg) <= remaining:
            yield seg
            remaining -= len(seg)
        else:
            yield seg[:remaining]
            break


def _apply_overlap(chunks: Iterable[str], overlap_size: int) -> Iterator[str]:
    """인접 청크가 서로 겹치게 만든다 — 경계에서 끊긴 문맥을 메운다.

    Args:
        chunks: 이미 만들어진 청크 스트림.
        overlap_size: 앞 청크에서 가져올 글자 수. **0 이하면 겹치지 않는다**(그대로 통과).

    Yields:
        앞부분이 겹쳐진 청크. **첫 청크는 그대로** 나온다(앞이 없다).
    """
    prev: str | None = None
    for chunk in chunks:
        if prev is None:
            prev = chunk
            yield chunk
            continue
        if overlap_size <= 0:
            yield chunk
        else:
            prefix = prev[-overlap_size:] if len(prev) > overlap_size else prev
            yield (prefix + "\n" + chunk).strip()
        prev = chunk


def _iter_document_chunks(
    path: Path,
    *,
    file_kind: str | None,
    encoding: str,
    chunk_size: int,
    overlap_size: int,
    max_input_chars: int,
) -> Iterator[str]:
    """``file_kind`` 에 맞는 형식별 파서로 원문을 세그먼트로 뽑아 청크 스트림으로 흘린다(공개 API의 실제 구현).

    형식이 무엇이든 공통 흐름은 같다: 파서가 낸 세그먼트 → 상한까지만 통과 → 묶기 → 겹치기.

    Args:
        path: 대상 파일.
        file_kind: 파일 종류. ⚠️ **실제 확장자와 어긋나면 막는다** — 엉뚱한 파서를 태우면
            깨진 텍스트가 조용히 임베딩된다.
        encoding: 읽을 인코딩.
        chunk_size: 청크 글자 한도.
        overlap_size: 겹칠 글자 수. 청크 크기 이상이면 무한 반복이 되므로 막는다.
        max_input_chars: 문서에서 읽을 누적 글자 상한.

    Yields:
        청크 문자열.

    Raises:
        ValueError: 겹침이 음수이거나 청크 크기 이상일 때, 종류와 확장자가 어긋날 때.
    """
    if overlap_size < 0:
        raise ValueError("overlap_size는 0 이상이어야 합니다.")
    if overlap_size >= chunk_size:
        raise ValueError("overlap_size는 chunk_size보다 작아야 합니다.")

    ext = path.suffix.lower()
    if file_kind is None:
        raise ValueError("file_kind는 필수입니다. txt/pdf/json/word/excel/powerpoint 중 하나를 전달하세요.")

    if file_kind in _OFFICE_FILE_KINDS:
        allowed = _OFFICE_EXTENSIONS[file_kind]
        if ext not in allowed:
            raise ValueError(
                f"file_kind={file_kind!r} 일 때 허용 확장자는 {sorted(allowed)} 입니다. 현재: {ext!r}"
            )
        if file_kind == OfficeKind.WORD.value:
            from docx import Document  # type: ignore

            doc = Document(str(path))
            lines = (p.text for p in doc.paragraphs)
            limited = _limit_segments(lines, max_input_chars)
            yield from _apply_overlap(_chunk_from_segments(limited, chunk_size), overlap_size)
            return
        if file_kind == OfficeKind.EXCEL.value:
            from openpyxl import load_workbook  # type: ignore

            wb = load_workbook(filename=str(path), read_only=True, data_only=True)
            try:
                def _rows() -> Iterator[str]:
                    """시트를 순서대로 훑어 한 줄씩 흘려보낸다(시트 제목 → 행들).

                    전부 메모리에 올리지 않으려고 제너레이터로 둔다 — 큰 표도 상한까지만
                    읽고 멈출 수 있다.
                    """
                    for sheet in wb.worksheets:
                        yield f"[시트] {sheet.title}"
                        for row in sheet.iter_rows(values_only=True):
                            cells = [str(v).strip() for v in row if v is not None and str(v).strip()]
                            if cells:
                                yield " | ".join(cells)

                limited = _limit_segments(_rows(), max_input_chars)
                yield from _apply_overlap(_chunk_from_segments(limited, chunk_size), overlap_size)
            finally:
                wb.close()
            return
        if file_kind == OfficeKind.POWERPOINT.value:
            from pptx import Presentation  # type: ignore

            prs = Presentation(str(path))

            def _slides() -> Iterator[str]:
                """슬라이드를 순서대로 훑어 텍스트를 흘려보낸다(번호 → 도형 텍스트)."""
                for i, slide in enumerate(prs.slides, start=1):
                    yield f"[슬라이드] {i}"
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            yield str(shape.text)

            limited = _limit_segments(_slides(), max_input_chars)
            yield from _apply_overlap(_chunk_from_segments(limited, chunk_size), overlap_size)
            return
        raise ValueError("지원하지 않는 오피스 file_kind 입니다.")

    if file_kind in _TEXT_FILE_KINDS:
        if file_kind == MediaKind.TEXT.value and ext != ".txt":
            raise ValueError(f"file_kind='{MediaKind.TEXT.value}' 일 때는 .txt 파일이어야 합니다.")
        if file_kind == MediaKind.JSON.value and ext != ".json":
            raise ValueError(f"file_kind='{MediaKind.JSON.value}' 일 때는 .json 파일이어야 합니다.")
        enc = _choose_encoding(path, encoding)
        with path.open("r", encoding=enc, errors="replace") as f:
            def _parts() -> Iterator[str]:
                """파일을 조각내어 읽는다 — 큰 파일을 통째로 메모리에 올리지 않기 위해서다."""
                while True:
                    part = f.read(4096)
                    if not part:
                        break
                    yield part

            limited = _limit_segments(_parts(), max_input_chars)
            yield from _apply_overlap(_chunk_from_segments(limited, chunk_size), overlap_size)
        return

    if file_kind == MediaKind.PDF.value:
        if ext != ".pdf":
            raise ValueError(f"file_kind='{MediaKind.PDF.value}' 일 때는 .pdf 파일이어야 합니다.")
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(str(path))
        pages = ((page.extract_text() or "") for page in reader.pages)
        limited = _limit_segments(pages, max_input_chars)
        yield from _apply_overlap(_chunk_from_segments(limited, chunk_size), overlap_size)
        return

    raise ValueError(
        "지원 확장자: .txt, .pdf, .json 또는 오피스(.docx/.xlsx/.pptx + file_kind=word|excel|powerpoint)."
    )


def iter_plain_text_chunks(
    text: str,
    *,
    chunk_size: int,
    overlap_size: int = 0,
    max_input_chars: int = _MAX_INPUT_CHARS,
) -> Iterator[str]:
    """파일 없이 문자열만 받아 문서와 **같은 규칙**으로 청크를 낸다.

    전사 텍스트처럼 파일로 떨어지지 않은 입력을 위한 경로다. 규칙이 문서 쪽과 같아야
    같은 내용이 어디서 들어왔는지에 따라 다르게 쪼개지지 않는다.

    Args:
        text: 쪼갤 원문. 비어 있거나 공백뿐이면 아무것도 내지 않는다.
        chunk_size: 청크 글자 한도.
        overlap_size: 겹칠 글자 수.
        max_input_chars: 읽을 누적 글자 상한. **0 이하면 아무것도 내지 않는다**.

    Yields:
        청크 문자열.

    Raises:
        ValueError: 겹침이 음수이거나 청크 크기 이상일 때(후자는 무한 반복이 된다).
    """
    if overlap_size < 0:
        raise ValueError("overlap_size는 0 이상이어야 합니다.")
    if overlap_size >= chunk_size:
        raise ValueError("overlap_size는 chunk_size보다 작아야 합니다.")
    if max_input_chars <= 0:
        return
    if text is None or not str(text).strip():
        return
    limited = _limit_segments(iter([str(text)]), max_input_chars)
    yield from _apply_overlap(_chunk_from_segments(limited, chunk_size), overlap_size)


# Public API (recommended)
normalize_file_kind = _normalize_file_kind
choose_encoding = _choose_encoding
iter_document_chunks = _iter_document_chunks
MAX_INPUT_CHARS = _MAX_INPUT_CHARS

__all__ = [
    "MAX_INPUT_CHARS",
    "choose_encoding",
    "iter_document_chunks",
    "iter_plain_text_chunks",
    "normalize_file_kind",
]

